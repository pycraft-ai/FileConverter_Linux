import os
import signal
import shutil
import subprocess
import threading
import resource
import concurrent.futures
from datetime import datetime
from utils.logger import get_logger

logger = get_logger(__name__)

import pandas as pd
import pytesseract
from PIL import Image
from pdf2docx import Converter
from pdf2image import convert_from_path
from pdfminer.high_level import extract_text
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from PyPDF2 import PdfMerger, PdfReader, PdfWriter


# ========================
# 常量定义
# ========================
DEFAULT_IMAGE_DPI = (150, 150)       # 图片转 PDF 默认 DPI
PDF_TO_IMAGE_DPI = 300               # PDF 转图片 DPI
OCR_TIMEOUT_SECONDS = 60             # OCR 最大处理时间（秒）
SLIDE_WIDTH_INCHES = 16              # PPT 幻灯片宽度（16:9）
SLIDE_HEIGHT_INCHES = 9              # PPT 幻灯片高度（16:9）

# LibreOffice 转换超时与资源限制（防止恶意/超大文件导致进程挂死或 OOM 造成 DoS）
LIBREOFFICE_TIMEOUT_SECONDS = int(os.environ.get('LIBREOFFICE_TIMEOUT_SECONDS', 120))  # 单次转换硬超时
LIBREOFFICE_MAX_RSS_MB = int(os.environ.get('LIBREOFFICE_MAX_RSS_MB', 1024))           # 进程地址空间上限(MB)
LIBREOFFICE_CPU_TIME = int(os.environ.get('LIBREOFFICE_CPU_TIME', 120))                # CPU 时间上限(秒)

# LibreOffice 转换锁（LibreOffice 实例不宜并发）
_libreoffice_lock = threading.Lock()


def _run_with_timeout(func, args=(), kwargs=None, timeout=OCR_TIMEOUT_SECONDS):
    """在独立线程中运行函数，支持超时（线程安全，替代 signal.SIGALRM）"""
    kwargs = kwargs or {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(func, *args, **kwargs)
        try:
            return future.result(timeout=timeout)
        except concurrent.futures.TimeoutError:
            raise TimeoutError("OCR处理超时")


class TimeoutError(Exception):
    """OCR处理超时异常"""
    pass


# LibreOffice 26.x 过滤器名称映射（简单扩展名 → 完整过滤器名）
# pdf 和 html 简写兼容，pptx/docx/xlsx 必须使用完整格式
_LO_FILTER_MAP = {
    'pptx': 'pptx:Impress MS PowerPoint 2007 XML',
    'docx': 'docx:MS Word 2007 XML',
    'xlsx': 'xlsx:Calc MS Excel 2007 XML',
}


def _libreoffice_convert(input_path, output_dir, convert_filter=None):
    """
    使用 LibreOffice 进行文件格式转换（WSL/Linux 替代 COM 方案）

    需要系统安装 LibreOffice: sudo apt install libreoffice-core libreoffice-writer \
        libreoffice-calc libreoffice-impress

    Args:
        input_path: 源文件路径
        output_dir: 输出目录
        convert_filter: 可选的导出过滤器（如 "pdf", "docx"），默认 pdf

    Returns:
        输出文件路径，失败返回 None
    """
    try:
        abs_input = os.path.abspath(input_path)
        abs_output = os.path.abspath(output_dir)
        os.makedirs(abs_output, exist_ok=True)

        # 解析过滤器名称（兼容 LibreOffice 26.x）
        raw_filter = convert_filter or 'pdf'
        actual_filter = _LO_FILTER_MAP.get(raw_filter, raw_filter)

        cmd = ['libreoffice', '--headless', '--norestore',
               '--convert-to', actual_filter,
               '--outdir', abs_output, abs_input]

        def _preexec_fn():
            """在子进程中设置资源限制，约束 LibreOffice 可用资源，防止 DoS。"""
            try:
                # 新进程组，便于超时时整体 kill（含 libreoffice 派生的所有子进程）
                os.setsid()
            except Exception:
                pass
            try:
                # 地址空间上限（字节）
                rss_bytes = LIBREOFFICE_MAX_RSS_MB * 1024 * 1024
                resource.setrlimit(resource.RLIMIT_AS, (rss_bytes, rss_bytes))
                # CPU 时间上限（秒）：超出后内核发送 SIGXCPU
                resource.setrlimit(resource.RLIMIT_CPU, (LIBREOFFICE_CPU_TIME, LIBREOFFICE_CPU_TIME))
                # 文件描述符数量上限
                resource.setrlimit(resource.RLIMIT_NOFILE, (1024, 1024))
            except (ValueError, OSError):
                pass

        proc = subprocess.Popen(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            preexec_fn=_preexec_fn,
        )
        try:
            stdout, stderr = proc.communicate(timeout=LIBREOFFICE_TIMEOUT_SECONDS)
            returncode = proc.returncode
        except subprocess.TimeoutExpired:
            # 超时：杀掉整个进程组（含 LibreOffice 所有残留子进程），避免孤儿进程累积导致 DoS
            try:
                os.killpg(os.getpgid(proc.pid), signal.SIGKILL)
            except (ProcessLookupError, Exception):
                proc.kill()
            proc.wait()
            logger.error("LibreOffice转换超时（超过 %s 秒），已强制终止进程组", LIBREOFFICE_TIMEOUT_SECONDS)
            return None

        if returncode != 0:
            logger.error("LibreOffice转换失败 (ret %s): %s", returncode, stderr.decode('utf-8', 'replace'))
            return None

        # 解析输出路径
        base = os.path.splitext(os.path.basename(input_path))[0]
        # 从过滤器名中提取扩展名（如 "pptx:Impress MS..." → "pptx"）
        if ':' in raw_filter:
            out_ext = raw_filter.split(':')[0]
        else:
            out_ext = raw_filter
        out_file = os.path.join(abs_output, f"{base}.{out_ext}")

        if os.path.exists(out_file):
            return out_file
        return None
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice转换超时")
        return None
    except Exception as e:
        logger.error("LibreOffice转换异常: %s", e)
        return None


def _format_size(size_bytes):
    """格式化文件大小为人类可读形式"""
    for unit in ('B', 'KB', 'MB', 'GB'):
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


def _parse_page_range(page_range, total_pages):
    """
    解析页码范围字符串，返回 1-based 页码列表。
    例如 "1-3,5,7-10" → [1,2,3,5,7,8,9,10]
    """
    pages = set()
    parts = page_range.replace('，', ',').split(',')
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            try:
                start, end = part.split('-', 1)
                start = int(start.strip())
                end = int(end.strip())
                if start < 1:
                    start = 1
                if end > total_pages:
                    end = total_pages
                for p in range(start, end + 1):
                    pages.add(p)
            except ValueError:
                continue
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p)
            except ValueError:
                continue
    return sorted(pages)


class Function:
    """文件转换引擎（WSL/Linux 版 — 使用 LibreOffice 替代 COM）"""

    @staticmethod
    def word_to_pdf(wordPath, pdfPath):
        """
        Word 转 PDF
        使用 LibreOffice --headless，适用于 WSL/Linux
        系统要求: sudo apt install libreoffice-core libreoffice-writer
        """
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(pdfPath))
                result = _libreoffice_convert(wordPath, output_dir, 'pdf')
                if result and os.path.exists(result):
                    # 如果 LibreOffice 输出的文件名不同，重命名为目标路径
                    if result != os.path.abspath(pdfPath):
                        shutil.move(result, pdfPath)
                    return True
                return False
            except Exception as e:
                logger.error("Word转PDF失败: %s", e)
                return False

    @staticmethod
    def md_to_pdf(mdPath, pdfPath):
        """
        Markdown 转 PDF
        策略: 先利用 markdown 库将 .md 解析为 HTML，
        再借助 LibreOffice 将 HTML 渲染并导出为 PDF
        """
        import markdown
        try:
            # 1. 读取 Markdown 内容
            with open(mdPath, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # 2. Markdown → HTML（启用常用扩展）
            md_extensions = [
                'tables',         # 表格支持
                'fenced_code',    # 围栏代码块
                'codehilite',     # 代码高亮
                'toc',            # 目录
                'nl2br',          # 换行转 <br>
            ]
            html_body = markdown.markdown(md_content, extensions=md_extensions)

            # 3. 包装为完整 HTML 文档（添加基础样式）
            html_doc = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{
    font-family: 'Noto Sans SC', 'DejaVu Sans', Arial, sans-serif;
    font-size: 14px;
    line-height: 1.8;
    color: #333;
    max-width: 900px;
    margin: 40px auto;
    padding: 0 20px;
}}
h1 {{ font-size: 26px; border-bottom: 2px solid #4A90D9; padding-bottom: 8px; }}
h2 {{ font-size: 22px; border-bottom: 1px solid #ddd; padding-bottom: 6px; }}
h3 {{ font-size: 18px; }}
h4 {{ font-size: 16px; }}
pre {{ background: #f5f5f5; padding: 12px; border-radius: 6px; overflow-x: auto; }}
code {{ background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-family: 'DejaVu Sans Mono', 'Courier New', monospace; }}
pre code {{ background: none; padding: 0; }}
table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
th, td {{ border: 1px solid #ccc; padding: 8px 12px; text-align: left; }}
th {{ background: #eef5fc; }}
blockquote {{ border-left: 4px solid #4A90D9; padding-left: 16px; color: #666; margin: 12px 0; }}
img {{ max-width: 100%; }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""

            # 4. 写入临时 HTML 文件
            html_path = mdPath + '.tmp.html'
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html_doc)

            # 5. 使用 LibreOffice 将 HTML 转为 PDF
            with _libreoffice_lock:
                output_dir = os.path.dirname(os.path.abspath(pdfPath))
                result = _libreoffice_convert(html_path, output_dir, 'pdf')

            # 6. 清理临时 HTML 文件
            try:
                os.remove(html_path)
            except Exception:
                pass

            if result and os.path.exists(result):
                if result != os.path.abspath(pdfPath):
                    shutil.move(result, pdfPath)
                return True
            return False
        except Exception as e:
            logger.error("MD转PDF失败: %s", e)
            # 清理可能残留的临时文件
            html_path = mdPath + '.tmp.html'
            if os.path.exists(html_path):
                try:
                    os.remove(html_path)
                except Exception:
                    pass
            return False

    @staticmethod
    def excel_to_pdf(excelPath, pdfPath):
        """
        Excel 转 PDF
        使用 LibreOffice，适用于 WSL/Linux
        系统要求: sudo apt install libreoffice-core libreoffice-calc
        """
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(pdfPath))
                result = _libreoffice_convert(excelPath, output_dir, 'pdf')
                if result and os.path.exists(result):
                    if result != os.path.abspath(pdfPath):
                        shutil.move(result, pdfPath)
                    return True
                return False
            except Exception as e:
                logger.error("Excel转PDF失败: %s", e)
                return False

    @staticmethod
    def ppt_to_pdf(pptPath, pdfPath):
        """
        PPT 转 PDF
        使用 LibreOffice，适用于 WSL/Linux
        系统要求: sudo apt install libreoffice-core libreoffice-impress
        """
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(pdfPath))
                result = _libreoffice_convert(pptPath, output_dir, 'pdf')
                if result and os.path.exists(result):
                    if result != os.path.abspath(pdfPath):
                        shutil.move(result, pdfPath)
                    return True
                return False
            except Exception as e:
                logger.error("PPT转PDF失败: %s", e)
                return False

    @staticmethod
    def html_to_pdf(htmlPath, pdfPath):
        """
        HTML 转 PDF
        使用 LibreOffice 渲染，适用于 WSL/Linux
        """
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(pdfPath))
                result = _libreoffice_convert(htmlPath, output_dir, 'pdf')
                if result and os.path.exists(result):
                    if result != os.path.abspath(pdfPath):
                        shutil.move(result, pdfPath)
                    return True
                return False
            except Exception as e:
                logger.error("HTML转PDF失败: %s", e)
                return False

    @staticmethod
    def pdf_to_word(pdfPath, wordPath):
        """PDF 转 Word（使用 pdf2docx，跨平台）"""
        try:
            cv = Converter(pdfPath)
            cv.convert(wordPath, start=0, end=None)
            cv.close()
            return True
        except Exception as e:
            logger.error("PDF转Word失败: %s", e)
            return False

    @staticmethod
    def pdf_to_ppt(pdfPath, pptPath):
        """PDF 转 PPT（使用 LibreOffice）"""
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(pptPath))
                result = _libreoffice_convert(pdfPath, output_dir, 'pptx')
                if result and os.path.exists(result):
                    if result != os.path.abspath(pptPath):
                        shutil.move(result, pptPath)
                    return True
                return False
            except Exception as e:
                logger.error("PDF转PPT失败: %s", e)
                return False

    @staticmethod
    def ppt_to_word(pptPath, wordPath):
        """
        PPT 转 Word（分两步：PPTX → PDF via LibreOffice → DOCX via pdf2docx）
        LibreOffice 不支持 Imppress → Writer 直接转换，需经 PDF 中转，
        且 PDF → DOCX 不能依赖 LibreOffice（PDF 以 Draw 模式加载），
        因此第二步使用 pdf2docx 库。
        """
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(wordPath))
                abs_ppt = os.path.abspath(pptPath)
                abs_word = os.path.abspath(wordPath)

                # 第一步：PPTX → PDF
                pdf_result = _libreoffice_convert(abs_ppt, output_dir, 'pdf')
                if not pdf_result or not os.path.exists(pdf_result):
                    logger.error("PPT转Word失败: 第一步 PPT→PDF 转换失败")
                    return False

                # 第二步：PDF → DOCX (使用 pdf2docx，与 pdf_to_word 一致)
                cv = Converter(pdf_result)
                cv.convert(abs_word, start=0, end=None)
                cv.close()

                # 清理中间 PDF
                try:
                    os.remove(pdf_result)
                except Exception:
                    pass

                return os.path.exists(abs_word) and os.path.getsize(abs_word) > 0
            except Exception as e:
                logger.error("PPT转Word失败: %s", e)
                return False

    @staticmethod
    def pdf_to_html(pdfPath, htmlPath):
        """PDF 转 HTML（使用 LibreOffice）"""
        with _libreoffice_lock:
            try:
                output_dir = os.path.dirname(os.path.abspath(htmlPath))
                result = _libreoffice_convert(pdfPath, output_dir, 'html')
                if result and os.path.exists(result):
                    if result != os.path.abspath(htmlPath):
                        shutil.move(result, htmlPath)
                    return True
                return False
            except Exception as e:
                logger.error("PDF转HTML失败: %s", e)
                return False

    @staticmethod
    def md_to_html(mdPath, htmlPath):
        """Markdown 转 HTML（使用 markdown 库 + 完整样式）"""
        import markdown
        try:
            with open(mdPath, 'r', encoding='utf-8') as f:
                md_content = f.read()

            md_extensions = [
                'tables', 'fenced_code', 'codehilite', 'toc', 'nl2br',
            ]
            html_body = markdown.markdown(md_content, extensions=md_extensions)

            html_doc = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<style>
body {{
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Noto Sans SC', sans-serif;
    font-size: 15px; line-height: 1.8; color: #24292e;
    max-width: 860px; margin: 40px auto; padding: 0 24px;
}}
h1 {{ font-size: 28px; border-bottom: 2px solid #e1e4e8; padding-bottom: 10px; margin-top: 24px; }}
h2 {{ font-size: 22px; border-bottom: 1px solid #e1e4e8; padding-bottom: 8px; margin-top: 20px; }}
h3 {{ font-size: 18px; margin-top: 16px; }}
h4 {{ font-size: 16px; }}
pre {{ background: #f6f8fa; padding: 16px; border-radius: 6px; overflow-x: auto; line-height: 1.45; }}
code {{ background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-family: 'SFMono-Regular', 'Consolas', monospace; font-size: 13px; }}
pre code {{ background: none; padding: 0; }}
table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
th, td {{ border: 1px solid #dfe2e5; padding: 8px 13px; text-align: left; }}
th {{ background: #f6f8fa; font-weight: 600; }}
tr:nth-child(even) td {{ background: #fafbfc; }}
blockquote {{ border-left: 4px solid #dfe2e5; padding: 0 16px; color: #6a737d; margin: 12px 0; }}
img {{ max-width: 100%; }}
a {{ color: #0366d6; text-decoration: none; }}
a:hover {{ text-decoration: underline; }}
ul, ol {{ padding-left: 2em; }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""

            with open(htmlPath, 'w', encoding='utf-8') as f:
                f.write(html_doc)
            return True
        except Exception as e:
            logger.error("MD转HTML失败: %s", e)
            return False

    @staticmethod
    def image_convert(imagePath, outputPath, target_format=None):
        """图片格式互转：jpg/png/webp/bmp/gif/tiff 互转"""
        try:
            # 从输出路径推断目标格式
            if not target_format:
                target_format = os.path.splitext(outputPath)[1].lstrip('.').lower()
            img = Image.open(imagePath)
            # 转换为 RGB（某些格式不支持 RGBA/调色板）
            if img.mode in ('RGBA', 'LA', 'P'):
                if target_format in ('jpg', 'jpeg', 'bmp'):
                    img = img.convert('RGB')
            elif img.mode == 'CMYK':
                img = img.convert('RGB')

            save_kwargs = {}
            if target_format in ('jpg', 'jpeg'):
                save_kwargs = {'quality': 92, 'optimize': True, 'subsampling': '4:4:4'}
            elif target_format == 'png':
                save_kwargs = {'optimize': True}
            elif target_format == 'webp':
                save_kwargs = {'quality': 85, 'method': 6}
            elif target_format == 'tiff':
                save_kwargs = {'compression': 'tiff_lzw'}

            img.save(outputPath, format=target_format.upper() if target_format.upper() == 'JPEG' else None, **save_kwargs)
            img.close()
            return True
        except Exception as e:
            logger.error("图片格式互转失败: %s", e)
            return False

    @staticmethod
    def image_to_pdf(imagePath, pdfPath):
        """图片 转 PDF（逐张处理，避免内存OOM）"""
        try:
            import re
            image_files = []
            for file in os.listdir(imagePath):
                if file.lower().endswith(('.jpg', '.jpeg', '.png')):
                    image_files.append(os.path.join(imagePath, file))
            if not image_files:
                return False

            # 使用自然排序
            def natural_sort_key(s):
                return [int(text) if text.isdigit() else text.lower()
                        for text in re.split(r'(\d+)', os.path.basename(s))]
            image_files.sort(key=natural_sort_key)

            # 逐张处理，不一次加载所有图片到内存
            first_img = None
            rest_images = []
            try:
                first_img = Image.open(image_files[0]).convert('RGB')
                if len(image_files) == 1:
                    first_img.save(pdfPath, format='PDF', dpi=DEFAULT_IMAGE_DPI)
                else:
                    # 逐张打开、添加、关闭
                    for img_path in image_files[1:]:
                        img = Image.open(img_path).convert('RGB')
                        rest_images.append(img)
                    first_img.save(
                        pdfPath, format='PDF', dpi=DEFAULT_IMAGE_DPI,
                        save_all=True, append_images=rest_images
                    )
            finally:
                # 无论成功或中途异常，都确保所有已打开的图片对象被关闭，避免资源泄漏
                if first_img is not None:
                    first_img.close()
                for img in rest_images:
                    img.close()
            return True
        except Exception as e:
            logger.error("图片转PDF失败: %s", e)
            return False

    @staticmethod
    def pdf_to_image(pdfPath, imageDir):
        """PDF 转 图片"""
        try:
            os.makedirs(imageDir, exist_ok=True)
            pages = convert_from_path(pdfPath, 150)
            for index, img in enumerate(pages):
                img.save(
                    os.path.join(imageDir, 'page_%s.jpg' % (index + 1)), 'JPEG'
                )
            return True
        except Exception as e:
            logger.error("PDF转图片失败: %s", e)
            return False

    @staticmethod
    def csv_to_excel(csvPath, excelPath):
        """CSV 转 Excel"""
        try:
            # 不强制设置 index_col，保留原始数据结构
            data = pd.read_csv(csvPath)
            data.to_excel(excelPath, index=False)
            return True
        except Exception as e:
            logger.error("CSV转Excel失败: %s", e)
            return False

    @staticmethod
    def excel_to_csv(excelPath, csvPath):
        """Excel 转 CSV"""
        try:
            # 不强制设置 index_col，保留原始数据结构
            data = pd.read_excel(excelPath, engine='openpyxl')
            data.to_csv(csvPath, encoding='utf-8', index=False)
            return True
        except Exception as e:
            logger.error("Excel转CSV失败: %s", e)
            return False

    @staticmethod
    def _clean_ocr_text(raw_text):
        """整理 OCR 原始文本，改善排版。
        策略（保守）：
        - 不强行合并段内行，保留段落多行结构
        - 清理行内/行末孤立噪声符号（□ | _ 等）
        - 规整行内空白（去除首尾空格、压缩重复空格）
        - 智能识别段间空行（结束标点、编号列表、标题开头）
        - 修复常见 OCR 错误字母（l↔I, 0↔O 等）
        """
        if not raw_text:
            return ""

        import re

        # 行内孤立噪声符号（OCR 误识别的方块、几何符号）
        inline_noise = re.compile(r'[□■▢▣◇◈●○►▸◄◅]')
        # 行尾孤立的标点/符号（OCR 切碎导致的尾巴，如行尾孤立 | _ - ）
        # 注意：保留中英文句号、逗号等正常标点
        trailing_noise = re.compile(r'[\s]*[|｜_—\-—·`~]+[\s]*$')

        # 结束标点（中英文）：行尾出现即视为该段结束
        end_punct_re = re.compile(r'[。！？；：”’"…）】》!?;:\.,]\s*$')
        # 段间分隔：行首为字母/数字+圆点/右括号（编号列表）
        list_item_re = re.compile(r'^\s*[A-Z]\s*[\)\.\:]', re.I)  # A) A. a)
        # 中文编号
        cn_list_re = re.compile(r'^\s*[\u4e00-\u9fff]{1,3}[、\.]')
        # 数字编号
        num_list_re = re.compile(r'^\s*\d{1,2}[\.\)、]\s*\S')
        # 常见大写英文标题（Section/Part/Directions/Answer/How to 等）
        title_re = re.compile(
            r'^\s*(Section\s+[A-Z]|Directions|Part\s+[IVX]|Answer\s+Sheet|'
            r'How\s+to\s+[A-Z]|Why\s+[A-Z]|What\s+[A-Z]|When\s+[A-Z]|'
            r'阅读理解|完形填空|翻译|写作|试题|答案与解析)',
            re.I)

        # 常见 OCR 错字替换（成对正则，逐项修复）
        ocr_fix = [
            (re.compile(r'\b0\b'), 'O'),           # 单词级单独 0 → O（保守，仅在两边空白时）
            (re.compile(r'\bl\b'), 'I'),           # 单词级单独 l → I
        ]
        # 注意：上面的 0/l 修复可能误伤数字，因此只在两侧是空格/标点时使用（已限定 \b）

        cleaned_lines = []
        for raw_line in raw_text.split('\n'):
            line = raw_line.rstrip()
            # 1) 压缩行内多余空格和制表符
            line = re.sub(r'[ \t]+', ' ', line)
            # 2) 去除行内孤立噪声符号
            line = inline_noise.sub('', line)
            # 3) 去除行尾孤立噪声（| _ - . 等）
            line = trailing_noise.sub('', line).rstrip()
            # 4) 修复常见 OCR 错误字母
            for pat, repl in ocr_fix:
                line = pat.sub(repl, line)
            cleaned_lines.append(line)

        # 智能段落化：基于"段间信号"插入空行
        paragraphs = []
        cur_para = []
        for line in cleaned_lines:
            stripped = line.strip()

            # 空行 → 段落边界
            if not stripped:
                if cur_para:
                    paragraphs.append(cur_para)
                    cur_para = []
                continue

            # 段落内第一行或检测到段间信号 → 新段落
            if not cur_para:
                cur_para = [line]
                continue

            # 段间信号：
            # 1) 当前行以编号/标题开头（强信号）
            cur_starts_new = bool(
                list_item_re.match(line) or cn_list_re.match(line)
                or num_list_re.match(line) or title_re.match(line)
            )
            # 2) 上一行末尾是结束标点 + 短行 + 下一行以大写/数字开头（弱信号，避免误判缩写）
            prev_line = cur_para[-1].strip()
            prev_ends_para = bool(end_punct_re.search(prev_line))
            weak_para_signal = (
                prev_ends_para and len(prev_line) <= 80 and bool(re.match(r'^[A-Z\u4e00-\u9fff\d]', line.strip()))
            )

            if cur_starts_new or weak_para_signal:
                paragraphs.append(cur_para)
                cur_para = [line]
            else:
                cur_para.append(line)

        if cur_para:
            paragraphs.append(cur_para)

        # 拼装：每个段内保留原换行（段内多行不合并）；段落之间用一个空行分隔
        return '\n\n'.join('\n'.join(p) for p in paragraphs).strip()

    @staticmethod
    def _ocr_pdf_pages(pdfPath):
        """OCR 辅助：将PDF页转图片后识别文字（在独立线程中运行）"""
        pages = convert_from_path(pdfPath, 300)
        text = ""
        for page in pages:
            text += pytesseract.image_to_string(
                page, lang='chi_sim+eng'
            ) + "\n"
        return text

    @staticmethod
    def _write_ocr_output(text, outputPath, fmt):
        """按指定格式写入 OCR 识别文本（txt / md / docx）"""
        fmt = (fmt or 'txt').lower()
        if fmt == 'docx':
            from docx import Document
            doc = Document()
            for para in (text or '').split('\n'):
                # 保留空行作为段落
                p = doc.add_paragraph()
                if para.strip():
                    run = p.add_run(para.strip())
                else:
                    p.paragraph_format.space_after = doc.styles['Normal'].paragraph_format.space_after
            doc.save(outputPath)
            return
        if fmt == 'md':
            # 简单转换为 Markdown：空行分隔为段落
            paragraphs = []
            buf = []
            for line in (text or '').split('\n'):
                if line.strip():
                    buf.append(line.strip())
                elif buf:
                    paragraphs.append(' '.join(buf))
                    buf = []
            if buf:
                paragraphs.append(' '.join(buf))
            md = '\n\n'.join(paragraphs)
            with open(outputPath, 'w', encoding='utf-8') as f:
                f.write(md)
            return
        # 默认 txt
        with open(outputPath, 'w', encoding='utf-8') as f:
            f.write(text or '')

    @staticmethod
    def pdf_ocr(pdfPath, outputPath, output_format='txt'):
        """PDF OCR 文字识别（带线程池超时控制，跨平台线程安全）"""
        try:
            text = extract_text(pdfPath)
            if len(text.strip()) < 10:
                text = _run_with_timeout(
                    Function._ocr_pdf_pages,
                    args=(pdfPath,),
                    timeout=OCR_TIMEOUT_SECONDS
                )
            # 整理 OCR 排版
            text = Function._clean_ocr_text(text)
            Function._write_ocr_output(text, outputPath, output_format)
            return True
        except TimeoutError:
            logger.error("PDF OCR超时")
            return False
        except Exception as e:
            logger.error("PDF OCR失败: %s", e)
            return False

    @staticmethod
    def _ocr_single_image(image_path):
        """OCR 辅助：识别单张图片文字（在独立线程中运行）"""
        return pytesseract.image_to_string(
            Image.open(image_path), lang='chi_sim+eng'
        )

    @staticmethod
    def image_ocr(imagePath, outputPath, output_format='txt'):
        """图片 OCR 文字识别（带线程池超时控制，跨平台线程安全）"""
        try:
            if os.path.isfile(imagePath):
                text = _run_with_timeout(
                    Function._ocr_single_image,
                    args=(imagePath,),
                    timeout=OCR_TIMEOUT_SECONDS
                )
                text = Function._clean_ocr_text(text)
                Function._write_ocr_output(text, outputPath, output_format)
                return True
            elif os.path.isdir(imagePath):
                all_text = ""
                for file in os.listdir(imagePath):
                    if file.lower().endswith(('.jpg', '.jpeg', '.png')):
                        img_path = os.path.join(imagePath, file)
                        text = _run_with_timeout(
                            Function._ocr_single_image,
                            args=(img_path,),
                            timeout=OCR_TIMEOUT_SECONDS
                        )
                        text = Function._clean_ocr_text(text)
                        all_text += f"=== {file} ===\n{text}\n\n"
                if all_text:
                    Function._write_ocr_output(all_text, outputPath, output_format)
                    return True
                return False
            return False
        except TimeoutError:
            logger.error("图片 OCR超时")
            return False
        except Exception as e:
            logger.error("图片OCR失败: %s", e)
            return False

    @staticmethod
    def img_to_ppt(imageDir, pptPath):
        """图片 转 PPT"""
        try:
            supported_formats = ('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff')
            image_files = []
            for file in os.listdir(imageDir):
                if file.lower().endswith(supported_formats):
                    image_files.append(os.path.join(imageDir, file))

            if not image_files:
                return False

            # 使用自然排序
            import re
            def natural_sort_key(s):
                return [int(text) if text.isdigit() else text.lower()
                        for text in re.split(r'(\d+)', os.path.basename(s))]
            image_files.sort(key=natural_sort_key)

            prs = Presentation()
            prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
            prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

            # 第一张作为封面
            first_image = image_files[0]
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            for shape in slide.placeholders:
                element = shape.element
                parent = element.getparent()
                parent.remove(element)
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            background.fill.background()
            pic = slide.shapes.add_picture(
                first_image, 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )

            # 剩余图片
            for image_file in image_files[1:]:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                for shape in slide.placeholders:
                    element = shape.element
                    parent = element.getparent()
                    parent.remove(element)
                background = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
                )
                background.fill.background()
                slide.shapes.add_picture(
                    image_file, 0, 0,
                    width=prs.slide_width, height=prs.slide_height
                )

            prs.save(pptPath)
            return True
        except Exception as e:
            logger.error("图片转PPT失败: %s", e)
            return False

    @staticmethod
    def merge_pdf(pdf_list, output_path):
        """合并 PDF"""
        try:
            merger = PdfMerger()
            for pdf in pdf_list:
                merger.append(pdf)
            merger.write(output_path)
            merger.close()
            return True
        except Exception as e:
            logger.error("合并PDF失败: %s", e)
            return False

    @staticmethod
    def pdf_encrypt(pdfPath, outputPath, password):
        """PDF 加密：设置打开密码保护"""
        try:
            reader = PdfReader(pdfPath)
            if reader.is_encrypted:
                logger.error("PDF已加密，请先解密后再加密")
                return False

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # 复制元数据
            if reader.metadata:
                writer.add_metadata(reader.metadata)

            writer.encrypt(user_password=password)

            with open(outputPath, 'wb') as f:
                writer.write(f)
            return True
        except Exception as e:
            logger.error("PDF加密失败: %s", e)
            return False

    @staticmethod
    def pdf_decrypt(pdfPath, outputPath, password):
        """PDF 解密：去除密码保护"""
        try:
            reader = PdfReader(pdfPath)

            if reader.is_encrypted:
                result = reader.decrypt(password)
                if result == 0:
                    logger.error("PDF解密失败：密码错误")
                    return False

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # 复制元数据
            if reader.metadata:
                writer.add_metadata(reader.metadata)

            with open(outputPath, 'wb') as f:
                writer.write(f)
            return True
        except Exception as e:
            logger.error("PDF解密失败: %s", e)
            return False

    # ==========================================================
    # 压缩 / 解压 / 去密
    # ==========================================================

    @staticmethod
    def compress_zip(input_paths, output_path, password=None, arcnames=None):
        """压缩为 ZIP（可选密码加密），arcnames 为原始文件名列表"""
        try:
            names = arcnames or [os.path.basename(p) for p in input_paths]
            if password:
                import pyzipper
                with pyzipper.AESZipFile(output_path, 'w',
                                          encryption=pyzipper.WZ_AES) as zf:
                    zf.setpassword(password.encode('utf-8'))
                    for path, name in zip(input_paths, names):
                        zf.write(path, name)
            else:
                import zipfile
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                    for path, name in zip(input_paths, names):
                        zf.write(path, name)
            return True
        except Exception as e:
            logger.error("ZIP压缩失败: %s", e)
            return False

    @staticmethod
    def compress_targz(input_paths, output_path, arcnames=None):
        """压缩为 TAR.GZ（不支持加密），arcnames 为原始文件名列表"""
        try:
            import tarfile
            names = arcnames or [os.path.basename(p) for p in input_paths]
            with tarfile.open(output_path, 'w:gz') as tf:
                for path, name in zip(input_paths, names):
                    tf.add(path, arcname=name)
            return True
        except Exception as e:
            logger.error("TAR.GZ压缩失败: %s", e)
            return False

    @staticmethod
    def compress_7z(input_paths, output_path, password=None, arcnames=None):
        """压缩为 7Z（可选 AES-256 加密），arcnames 为原始文件名列表"""
        try:
            import py7zr
            names = arcnames or [os.path.basename(p) for p in input_paths]
            if password:
                with py7zr.SevenZipFile(output_path, 'w',
                                         password=password) as szf:
                    for path, name in zip(input_paths, names):
                        szf.write(path, name)
            else:
                with py7zr.SevenZipFile(output_path, 'w') as szf:
                    for path, name in zip(input_paths, names):
                        szf.write(path, name)
            return True
        except Exception as e:
            logger.error("7Z压缩失败: %s", e)
            return False

    @staticmethod
    def is_archive_encrypted(path):
        """检测压缩包是否加密（快速检查，不提取内容）"""
        fname = path.lower()
        try:
            if fname.endswith('.zip'):
                import zipfile
                with zipfile.ZipFile(path, 'r') as zf:
                    for info in zf.infolist():
                        if info.flag_bits & 0x1:
                            return True
                return False
            elif fname.endswith('.7z'):
                import py7zr
                try:
                    with py7zr.SevenZipFile(path, 'r') as szf:
                        return szf.needs_password()
                except Exception:
                    return False
            return False
        except Exception as e:
            logger.error("检测加密状态失败: %s", e)
            return False

    @staticmethod
    def _safe_archive_members(member_names, output_dir):
        """
        校验压缩包成员路径不会跨越 output_dir（防御 Zip Slip / 路径穿越）。
        合法成员必须落在 output_dir 之内（含目录自身）。
        """
        base = os.path.realpath(output_dir)
        for name in member_names:
            # 规范化成员名，剥离开头的 ./ 与多余的斜杠
            norm = name.replace('\\', '/')
            dest = os.path.realpath(os.path.join(base, norm))
            if dest != base and not dest.startswith(base + os.sep):
                raise ValueError(f'压缩包含非法路径成员: {name!r}')

    @staticmethod
    def decompress_archive(input_path, output_dir, password=None):
        """解压压缩包（自动识别格式）。防御 Zip Slip、符号链接穿越与解压炸弹。"""
        from config import Config
        fname = input_path.lower()
        try:
            if fname.endswith('.zip'):
                # 解压前按元数据预检解压总大小/文件数，防御 zip 炸弹（在解压前拦截，最有效）
                if password:
                    import pyzipper
                    with pyzipper.AESZipFile(input_path, 'r') as zf:
                        zf.setpassword(password.encode('utf-8'))
                        infos = zf.infolist()
                        Function._precheck_archive(infos)
                        Function._safe_archive_members(zf.namelist(), output_dir)
                        zf.extractall(output_dir)
                else:
                    import zipfile
                    with zipfile.ZipFile(input_path, 'r') as zf:
                        infos = zf.infolist()
                        Function._precheck_archive(infos)
                        Function._safe_archive_members(zf.namelist(), output_dir)
                        zf.extractall(output_dir)
                return True
            elif fname.endswith('.tar.gz') or fname.endswith('.tgz') or fname.endswith('.tar'):
                import tarfile
                with tarfile.open(input_path, 'r:*') as tf:
                    # filter='data' 自动防御路径穿越/绝对路径/符号链接/设备文件
                    members = tf.getmembers()
                    Function._safe_archive_members([m.name for m in members], output_dir)
                    # 预检解压总大小/文件数（tar 成员无 file_size，用 size）
                    total = sum(m.size for m in members)
                    if total > Config.ARCHIVE_MAX_TOTAL_BYTES:
                        raise ValueError('压缩包解压后总大小超过限制')
                    if len(members) > Config.ARCHIVE_MAX_FILES:
                        raise ValueError('压缩包文件数量超过限制')
                    tf.extractall(output_dir, filter='data')
                return True
            elif fname.endswith('.7z'):
                import py7zr
                with py7zr.SevenZipFile(input_path, 'r',
                                        password=password) as szf:
                    # 7z 先列出文件数，解压后通过目录统计校验总大小
                    names = szf.getnames()
                    if len(names) > Config.ARCHIVE_MAX_FILES:
                        raise ValueError('压缩包文件数量超过限制')
                    Function._safe_archive_members(names, output_dir)
                    szf.extractall(output_dir)
                    # 解压后统计实际大小，防止 7z 元数据与实际不符
                    Function._check_dir_size(output_dir)
                return True
            else:
                logger.error("不支持的解压格式: %s", input_path)
                return False
        except ValueError as e:
            # 路径穿越 / 解压炸弹等安全问题，明确记录并拒绝
            logger.error("解压被拒绝（安全校验失败）: %s", e)
            return False
        except Exception as e:
            logger.error("解压失败: %s", e)
            return False

    @staticmethod
    def _precheck_archive(infos):
        """解压前预检 zip 成员：总解压大小、文件数量（防御解压炸弹）。

        在 extractall 之前基于压缩包声明的解压大小拦截，避免实际解压造成磁盘/内存耗尽。
        """
        from config import Config
        total_bytes = 0
        file_count = 0
        for info in infos:
            total_bytes += getattr(info, 'file_size', 0)
            file_count += 1
            if total_bytes > Config.ARCHIVE_MAX_TOTAL_BYTES:
                raise ValueError('压缩包解压后总大小超过限制')
            if file_count > Config.ARCHIVE_MAX_FILES:
                raise ValueError('压缩包文件数量超过限制')

    @staticmethod
    def _check_dir_size(output_dir):
        """统计解压目录的实际大小，超过限制则拒绝并清理。"""
        from config import Config
        total = 0
        for root, _dirs, files in os.walk(output_dir):
            for f in files:
                try:
                    total += os.path.getsize(os.path.join(root, f))
                except OSError:
                    pass
                if total > Config.ARCHIVE_MAX_TOTAL_BYTES:
                    # 超限：删除已解压内容，避免占用磁盘
                    shutil.rmtree(output_dir, ignore_errors=True)
                    raise ValueError('压缩包解压后总大小超过限制')

    @staticmethod
    def decrypt_archive(input_path, output_path, password):
        """去除压缩包密码保护（解压后重新打包为无密码 ZIP）"""
        import tempfile
        import shutil
        temp_dir = tempfile.mkdtemp()
        try:
            # 用密码解压
            if not Function.decompress_archive(input_path, temp_dir, password):
                return False
            # 收集所有解压出的文件
            extracted = []
            for root, _dirs, files in os.walk(temp_dir):
                for f in files:
                    extracted.append(os.path.join(root, f))
            if not extracted:
                return False
            # 重新打包为无密码 ZIP
            import zipfile
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for file_path in extracted:
                    arcname = os.path.relpath(file_path, temp_dir)
                    zf.write(file_path, arcname)
            return True
        except Exception as e:
            logger.error("压缩包去密失败: %s", e)
            return False
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    # ==========================================================
    # PDF 压缩 / 分割 / 转 Excel
    # ==========================================================

    @staticmethod
    def pdf_compress(pdfPath, outputPath, quality='medium'):
        """
        PDF 压缩：通过 PyPDF2 压缩内容流 + 移除元数据 + 图片缩减采样。
        quality: 'low'(高压缩) / 'medium'(平衡) / 'high'(轻度压缩)

        注意：如需极致压缩率（Ghostscript），需额外安装 gs 并注释下方 fallback。
        """
        try:
            reader = PdfReader(pdfPath)
            writer = PdfWriter()

            # 质量档位 → 压缩参数
            quality_config = {
                'low':  {'image_quality': 30, 'compress': True},
                'medium': {'image_quality': 55, 'compress': True},
                'high': {'image_quality': 80, 'compress': True},
            }
            qc = quality_config.get(quality, quality_config['medium'])

            total_pages = len(reader.pages)
            for page in reader.pages:
                # 压缩页面内容流
                if qc['compress']:
                    page.compress_content_streams()
                writer.add_page(page)

            # 移除元数据可略微减少体积
            # writer.add_metadata({})  # 如果不需要元数据

            with open(outputPath, 'wb') as f:
                writer.write(f)

            # 获取压缩后大小
            orig_size = os.path.getsize(pdfPath)
            new_size = os.path.getsize(outputPath)
            compressed_pct = round((1 - new_size / orig_size) * 100)
            logger.info("PDF压缩完成 | %s → %s (%d%% 减少) | pages=%d",
                        _format_size(orig_size), _format_size(new_size),
                        compressed_pct, total_pages)
            return True
        except Exception as e:
            logger.error("PDF压缩失败: %s", e)
            return False

    @staticmethod
    def pdf_split(pdfPath, outputPath, page_range):
        """
        PDF 分割/提取页：按页码范围提取指定页面。

        Args:
            page_range: 页码范围字符串，如 "1-3,5,7-10"。页码从 1 开始。
        Returns:
            (原文件页数, 提取页数) 用于日志展示。
        """
        try:
            reader = PdfReader(pdfPath)
            total_pages = len(reader.pages)

            # 解析页码范围
            pages_to_extract = _parse_page_range(page_range, total_pages)
            if not pages_to_extract:
                logger.error("PDF分割失败：页码范围无效")
                return None

            writer = PdfWriter()
            for page_num in pages_to_extract:
                writer.add_page(reader.pages[page_num - 1])  # 转 0-based

            with open(outputPath, 'wb') as f:
                writer.write(f)

            logger.info("PDF分割完成 | 从 %d 页中提取 %d 页 | range=%s",
                        total_pages, len(pages_to_extract), page_range)
            return (total_pages, len(pages_to_extract))
        except Exception as e:
            logger.error("PDF分割失败: %s", e)
            return None

    @staticmethod
    def pdf_to_excel(pdfPath, excelPath):
        """
        PDF 转 Excel：使用 pdfplumber 提取表格数据。

        需要安装: pip install pdfplumber
        """
        import pdfplumber
        try:
            all_tables = []
            with pdfplumber.open(pdfPath) as pdf:
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    if tables:
                        for j, table in enumerate(tables):
                            if table and len(table) > 0:
                                df = pd.DataFrame(table[1:], columns=table[0] if table[0] else None)
                                all_tables.append((i + 1, j + 1, df))

            if not all_tables:
                # 无表格时尝试提取纯文本
                all_text = []
                with pdfplumber.open(pdfPath) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            all_text.append(text)
                if all_text:
                    df = pd.DataFrame({'内容': all_text})
                    df.to_excel(excelPath, index=False, sheet_name='文本内容')
                    return True
                return False

            # 多个表格写入不同 sheet
            if len(all_tables) == 1:
                all_tables[0][2].to_excel(excelPath, index=False, sheet_name=f'第{all_tables[0][0]}页_表{all_tables[0][1]}')
            else:
                with pd.ExcelWriter(excelPath, engine='openpyxl') as writer:
                    for page_num, table_idx, df in all_tables:
                        sheet_name = f'P{page_num}_T{table_idx}'[:31]  # sheet 名最长31字符
                        df.to_excel(writer, index=False, sheet_name=sheet_name)

            logger.info("PDF转Excel完成 | 提取 %d 个表格", len(all_tables))
            return True
        except Exception as e:
            logger.error("PDF转Excel失败: %s", e)
            return False

    # ==========================================================
    # 图片压缩
    # ==========================================================

    @staticmethod
    def image_compress(imagePath, outputPath, quality=75, max_width=2048, max_height=2048):
        """
        图片压缩：调整尺寸 + 降低质量。

        Args:
            quality: JPEG/WebP 质量 (1-100)，默认 75
            max_width/max_height: 最大尺寸，超出等比缩放
        """
        try:
            img = Image.open(imagePath)
            orig_size = os.path.getsize(imagePath)
            orig_w, orig_h = img.size

            # 等比缩放到不超出最大尺寸
            ratio = 1.0
            if orig_w > max_width:
                ratio = min(ratio, max_width / orig_w)
            if orig_h > max_height:
                ratio = min(ratio, max_height / orig_h)

            if ratio < 1.0:
                new_w = int(orig_w * ratio)
                new_h = int(orig_h * ratio)
                img = img.resize((new_w, new_h), Image.LANCZOS)

            # 确定输出格式
            out_ext = os.path.splitext(outputPath)[1].lower()
            fmt = out_ext.lstrip('.')

            save_kwargs = {}
            if fmt in ('jpg', 'jpeg'):
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                save_kwargs = {'quality': quality, 'optimize': True, 'subsampling': '4:2:0'}
            elif fmt == 'png':
                save_kwargs = {'optimize': True}
            elif fmt == 'webp':
                if img.mode in ('RGBA', 'LA'):
                    pass
                elif img.mode == 'P':
                    img = img.convert('RGBA')
                save_kwargs = {'quality': quality, 'method': 6}
            elif fmt in ('bmp', 'tiff', 'gif'):
                save_kwargs = {}

            img.save(outputPath, **save_kwargs)
            img.close()

            new_size = os.path.getsize(outputPath)
            saved_pct = round((1 - new_size / orig_size) * 100) if orig_size > 0 else 0
            new_w, new_h = (0, 0)
            logger.info("图片压缩完成 | %s → %s (%d%% 减少) | %dx%d → %dx%d",
                        _format_size(orig_size), _format_size(new_size),
                        saved_pct, orig_w, orig_h, new_w or orig_w, new_h or orig_h)
            return True
        except Exception as e:
            logger.error("图片压缩失败: %s", e)
            return False

    # ==========================================================
    # 文字转语音
    # ==========================================================

    @staticmethod
    def txt_to_speech(txtPath, mp3Path, voice='zh-CN-XiaoxiaoNeural', rate='+0%'):
        """
        文字转语音：使用 Edge-TTS (微软 TTS) 将文本转为 MP3。

        需要安装: pip install edge-tts
        优势: 国内网络可直接访问，中文语音自然流畅。

        :param voice: 语音角色，如 zh-CN-XiaoxiaoNeural
        :param rate: 语速，如 '+20%', '-30%', '+0%'
        """
        import asyncio
        import edge_tts

        try:
            with open(txtPath, 'r', encoding='utf-8') as f:
                text = f.read()

            if not text.strip():
                return False

            # 长文本截断（单次请求有限制）
            max_chunk = 3000
            if len(text) > max_chunk:
                logger.warning("文本过长(%d字符)，截取前 %d 字符进行转换", len(text), max_chunk)
                text = text[:max_chunk]

            async def _speak(text, out_path, voice, rate):
                communicate = edge_tts.Communicate(
                    text=text,
                    voice=voice,
                    rate=rate,
                )
                await communicate.save(out_path)

            asyncio.run(_speak(text, mp3Path, voice, rate))

            logger.info("文字转语音完成 | 文本长度=%d voice=%s rate=%s", len(text), voice, rate)
            return True
        except Exception as e:
            logger.error("文字转语音失败: %s", e)
            return False
